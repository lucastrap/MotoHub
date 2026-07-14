from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
BLACK     = RGBColor(0x09, 0x09, 0x09)
DARK      = RGBColor(0x1a, 0x1a, 0x1a)
CARD      = RGBColor(0x22, 0x22, 0x22)
RED       = RGBColor(0xCC, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xAA, 0xAA)
LIGHTGRAY = RGBColor(0x77, 0x77, 0x77)
GREEN     = RGBColor(0x22, 0x88, 0x44)
BLUE      = RGBColor(0x22, 0x55, 0xAA)
ORANGE    = RGBColor(0xCC, 0x77, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

# ─── helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    try:
        run.font.name = "Barlow Condensed"
    except Exception:
        run.font.name = "Arial Narrow"
    return txb

def bg(slide):
    add_rect(slide, 0, 0, W, H, fill=BLACK)

def red_bar(slide):
    add_rect(slide, 0, 0, Inches(0.12), H, fill=RED)

def top_label(slide, text):
    add_rect(slide, Inches(0.18), 0, W, Inches(0.38), fill=RGBColor(0x18,0x18,0x18))
    add_text(slide, text, Inches(0.22), Inches(0.06), Inches(12), Inches(0.32),
             size=9, bold=True, color=RED)

def slide_title(slide, text, y=Inches(0.45)):
    add_text(slide, text, Inches(0.22), y, Inches(12.8), Inches(1.2),
             size=44, bold=True, color=WHITE)

def card(slide, x, y, w, h, fill=CARD, border=None):
    add_rect(slide, x, y, w, h, fill=fill,
             line_color=border or RGBColor(0x44,0x44,0x44), line_width=Pt(0.75))

def card_header(slide, text, x, y, w, h=Inches(0.38), fill=RED):
    add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, text, x+Inches(0.1), y+Inches(0.04), w-Inches(0.2), h,
             size=11, bold=True, color=WHITE)

def bullet(slide, lines, x, y, w, h, size=11, color=GRAY):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        try:
            run.font.name = "Barlow Condensed"
        except Exception:
            run.font.name = "Arial Narrow"

def footer_note(slide, text):
    add_rect(slide, Inches(0.22), Inches(6.9), Inches(12.8), Inches(0.45),
             fill=RGBColor(0x18,0x18,0x18),
             line_color=RED, line_width=Pt(0.75))
    add_text(slide, text, Inches(0.35), Inches(6.92), Inches(12.5), Inches(0.42),
             size=10, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
red_bar(s)
add_rect(s, 0, 0, Inches(0.12), H, fill=RED)

# Big title
add_text(s, "MotoTrack", Inches(0.3), Inches(1.0), Inches(8), Inches(2.2),
         size=90, bold=True, color=WHITE)
add_text(s, "Application web de gestion de moto",
         Inches(0.3), Inches(3.1), Inches(8), Inches(0.6),
         size=22, color=GRAY)

# Tech chips
chips = ["Next.js 14", "TypeScript", "PostgreSQL", "Prisma", "Three.js"]
cx = Inches(0.3)
for chip in chips:
    add_rect(s, cx, Inches(3.9), Inches(1.9), Inches(0.45), fill=DARK,
             line_color=RGBColor(0x44,0x44,0x44))
    add_text(s, chip, cx+Inches(0.05), Inches(3.93), Inches(1.8), Inches(0.4),
             size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
    cx += Inches(2.0)

add_rect(s, Inches(0.3), Inches(4.5), Inches(9.5), Inches(0.02), fill=GRAY)
add_text(s, "Projet de Fin d'Études  •  Bloc 1 — Cadrage de projet  •  RNCP 39583",
         Inches(0.3), Inches(4.65), Inches(9), Inches(0.4), size=12, color=LIGHTGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — C1.1.1 PARTIES PRENANTES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.1.1 — CARTOGRAPHIE DES PARTIES PRENANTES")
slide_title(s, "Qui est impliqué dans le projet ?")

cols = [
    ("Commanditaire", RED,    ["MotoClub Alpes", "Marc Duvalier", "Vision produit", "Validation finale"]),
    ("Équipe projet", BLUE,   ["Luca S. — Lead dev", "Thomas R. — Backend", "Sarah M. — Frontend", "Karim B. — Design", "Julie F. — QA"]),
    ("Utilisateurs",  GREEN,  ["Motards 18-45 ans", "1 à 3 motos", "Suivi entretien", "Recherche pièces"]),
    ("Admin système", ORANGE, ["Infrastructure", "PostgreSQL Docker", "CI/CD GitHub Actions", "Monitoring /health"]),
    ("APIs externes", RGBColor(0x66,0x33,0xAA), ["NHTSA (modèles)", "Google News RSS", "eBay / LBC / Amazon", "wttr.in (météo)"]),
]
cw = Inches(2.4)
cx = Inches(0.22)
for title, color, items in cols:
    card(s, cx, Inches(1.55), cw, Inches(5.1))
    card_header(s, title, cx, Inches(1.55), cw, fill=color)
    bullet(s, items, cx+Inches(0.1), Inches(2.05), cw-Inches(0.2), Inches(4.5), size=11)
    cx += cw + Inches(0.1)

footer_note(s, "Utilisateurs cibles : motards 18-45 ans, propriétaires de 1 à 3 motos, suivant régulièrement leurs entretiens")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — C1.1.2 ANALYSE DE LA DEMANDE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.1.2 — ANALYSE DE LA DEMANDE")
slide_title(s, "Problématique & objectifs")

# Problem box
add_rect(s, Inches(0.22), Inches(1.6), Inches(12.8), Inches(0.9),
         fill=RGBColor(0x15,0x15,0x15), line_color=RED, line_width=Pt(1))
add_text(s, "Problématique client", Inches(0.35), Inches(1.62), Inches(12), Inches(0.3),
         size=11, bold=True, color=RED)
add_text(s,
    "Les motards n'ont aucun outil centralisé pour gérer leur garage, suivre les entretiens et accéder "
    "à l'information (actualités, pièces, météo) depuis une seule interface.",
    Inches(0.35), Inches(1.9), Inches(12.5), Inches(0.55), size=12, color=GRAY)

objectives = [
    ("01", "Centraliser",
     "Un seul espace pour gérer plusieurs motos (garage virtuel) avec données complètes : "
     "marque, modèle, immat SIV, kilométrage, achat"),
    ("02", "Tracer",
     "Historique des maintenances (vidange, pneus, freins, chaîne…) avec rappels et types "
     "d'entretien normalisés"),
    ("03", "Informer",
     "Actualités RSS, pièces détachées multi-sources (eBay, LBC, Amazon), météo locale pour "
     "préparer les sorties"),
]
cw = Inches(4.1)
cx = Inches(0.22)
for num, title, desc in objectives:
    card(s, cx, Inches(2.65), cw, Inches(4.1))
    add_text(s, num, cx+Inches(0.15), Inches(2.75), Inches(0.8), Inches(0.7),
             size=36, bold=True, color=RED)
    add_text(s, title, cx+Inches(0.15), Inches(3.35), cw-Inches(0.3), Inches(0.45),
             size=16, bold=True, color=WHITE)
    bullet(s, [desc], cx+Inches(0.15), Inches(3.85), cw-Inches(0.3), Inches(2.8), size=11)
    cx += cw + Inches(0.13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — C1.2.1 SWOT (CORRIGÉ)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.2.1 — OPPORTUNITÉS & MENACES (SWOT)")
slide_title(s, "Analyse SWOT du projet")

quadrants = [
    ("FORCES",        GREEN,  Inches(0.22),  Inches(1.6),
     ["Stack moderne Next.js 14 + TypeScript",
      "Auth JWT httpOnly — sécurisée",
      "Scène 3D différenciante (Ducati V4R)",
      "Validation stricte Zod côté client ET API",
      "31 tests unitaires Jest — CI/CD GitHub Actions"]),
    ("FAIBLESSES",    ORANGE, Inches(6.72),  Inches(1.6),
     ["Auth maison (pas OAuth — Google/Apple)",
      "Pas d'app mobile native (PWA uniquement)",
      "Dépendance flux Google News RSS",
      "Pas encore déployé en production cloud"]),
    ("OPPORTUNITÉS",  BLUE,   Inches(0.22),  Inches(4.35),
     ["Marché moto en croissance (+7% France 2024)",
      "API NHTSA gratuite & fiable",
      "12 marques chinoises émergentes intégrées",
      "Monétisation possible (abonnement premium)"]),
    ("MENACES",       RED,    Inches(6.72),  Inches(4.35),
     ["Google peut modifier ses flux RSS",
      "Concurrence (Mapit, BikeReg, CarnetMoto)",
      "Impact RGPD sur données utilisateurs",
      "Instabilité APIs tierces (eBay, Amazon)"]),
]
for title, color, x, y, items in quadrants:
    card(s, x, y, Inches(6.3), Inches(2.6))
    card_header(s, title, x, y, Inches(6.3), fill=color)
    bullet(s, items, x+Inches(0.15), y+Inches(0.5), Inches(6.0), Inches(2.0), size=11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — C1.2.2 FAISABILITÉ TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.2.2 — FAISABILITÉ TECHNIQUE & AUDIT")
slide_title(s, "Environnement technique")

cols = [
    ("Frontend",        BLUE,
     ["Next.js 14 App Router",
      "TypeScript strict",
      "Tailwind CSS + shadcn/ui",
      "React Hook Form + Zod",
      "Three.js / R3F (3D)",
      "Recharts (graphiques)"]),
    ("Backend / API",   RED,
     ["Next.js API Routes",
      "JWT Auth (httpOnly cookie)",
      "Prisma ORM + migrations",
      "NHTSA API (modèles, cache 24h)",
      "Google News RSS x4",
      "REST endpoints sécurisés"]),
    ("Infrastructure",  GREEN,
     ["PostgreSQL 15 (Docker)",
      "Adminer (port 8080)",
      "docker-compose.yml",
      "GitHub Actions CI/CD",
      "Winston logger + /api/health",
      "Prêt déploiement Vercel/Supabase"]),
]
cw = Inches(4.1)
cx = Inches(0.22)
for title, color, items in cols:
    card(s, cx, Inches(1.6), cw, Inches(5.1))
    card_header(s, title, cx, Inches(1.6), cw, fill=color)
    bullet(s, items, cx+Inches(0.15), Inches(2.1), cw-Inches(0.3), Inches(4.5), size=12)
    cx += cw + Inches(0.13)

footer_note(s, "Contraintes : hébergement cloud requis, < 10 000 utilisateurs estimés, délai 6 semaines, équipe 5 personnes")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — C1.5 ARCHITECTURE LOGICIELLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.5 — ARCHITECTURE LOGICIELLE")
slide_title(s, "Architecture Next.js App Router")

# Client layer
add_rect(s, Inches(0.22), Inches(1.55), Inches(12.8), Inches(1.5),
         fill=RGBColor(0x1a,0x22,0x3a), line_color=BLUE, line_width=Pt(0.75))
add_text(s, "COUCHE CLIENT — Navigateur", Inches(0.35), Inches(1.58), Inches(4), Inches(0.3),
         size=9, bold=True, color=BLUE)
pages = ["/\nLanding 3D", "/garage\nListe motos", "/dashboard\nTableau bord",
         "/maintenance\nHistorique", "/news\nActualités", "/pieces\nPièces", "/weather\nMétéo"]
px = Inches(0.35)
for page in pages:
    lines = page.split("\n")
    add_rect(s, px, Inches(1.88), Inches(1.68), Inches(0.95),
             fill=RGBColor(0x1e,0x30,0x55), line_color=BLUE, line_width=Pt(0.5))
    add_text(s, lines[0], px+Inches(0.05), Inches(1.9), Inches(1.58), Inches(0.35),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, lines[1], px+Inches(0.05), Inches(2.22), Inches(1.58), Inches(0.3),
             size=9, color=GRAY, align=PP_ALIGN.CENTER)
    px += Inches(1.75)

add_text(s, "▼  API Routes", Inches(0), Inches(3.1), W, Inches(0.3),
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

# API layer
add_rect(s, Inches(0.22), Inches(3.38), Inches(12.8), Inches(1.5),
         fill=RGBColor(0x2a,0x10,0x10), line_color=RED, line_width=Pt(0.75))
add_text(s, "COUCHE API — Next.js API Routes + JWT Auth", Inches(0.35), Inches(3.4), Inches(6), Inches(0.3),
         size=9, bold=True, color=RED)
endpoints = [
    "POST /api/auth/register", "POST /api/auth/login",    "GET·POST /api/motorcycles",
    "PATCH /api/motorcycles/[id]","GET·POST /api/maintenances","GET /api/news",
    "GET /api/motorcycle-models", "POST /api/auth/logout",  "GET /api/health",
]
ex = Inches(0.35)
ey = Inches(3.72)
for i, ep in enumerate(endpoints):
    if i == 5:
        ex = Inches(0.35)
        ey = Inches(4.15)
    add_rect(s, ex, ey, Inches(4.05), Inches(0.35),
             fill=RGBColor(0x3a,0x18,0x18), line_color=RGBColor(0x66,0x22,0x22), line_width=Pt(0.5))
    add_text(s, ep, ex+Inches(0.08), ey+Inches(0.04), Inches(3.9), Inches(0.3),
             size=9, color=RGBColor(0xFF,0x88,0x88))
    ex += Inches(4.2)

add_text(s, "▼  Prisma ORM", Inches(0), Inches(4.93), W, Inches(0.28),
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Data layer
add_rect(s, Inches(0.22), Inches(5.18), Inches(12.8), Inches(1.6),
         fill=RGBColor(0x0e,0x22,0x18), line_color=GREEN, line_width=Pt(0.75))
add_text(s, "COUCHE DONNÉES — PostgreSQL + Prisma", Inches(0.35), Inches(5.21), Inches(6), Inches(0.3),
         size=9, bold=True, color=GREEN)
models = [
    ("User",        "id · email · passwordHash · name"),
    ("Motorcycle",  "id · brand · model · year · mileage · isPrimary"),
    ("Maintenance", "id · type · date · mileage · cost"),
    ("Reminder",    "id · type · thresholdMileage · frequencyMonths"),
]
mx = Inches(0.35)
for name, fields in models:
    add_rect(s, mx, Inches(5.52), Inches(3.05), Inches(1.1),
             fill=RGBColor(0x12,0x2e,0x20), line_color=GREEN, line_width=Pt(0.5))
    add_text(s, name, mx+Inches(0.1), Inches(5.55), Inches(2.85), Inches(0.35),
             size=11, bold=True, color=GREEN)
    add_text(s, fields, mx+Inches(0.1), Inches(5.88), Inches(2.85), Inches(0.65),
             size=9, color=GRAY)
    mx += Inches(3.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — C1.2.3 RISQUES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.2.3 — CARTOGRAPHIE DES RISQUES")
slide_title(s, "Risques techniques & fonctionnels")

# Table header
add_rect(s, Inches(0.22), Inches(1.6), Inches(12.8), Inches(0.4), fill=RED)
for text, x, w in [
    ("Risque",       Inches(0.35), Inches(3.2)),
    ("Catégorie",    Inches(3.65), Inches(1.5)),
    ("Prob.",        Inches(5.22), Inches(1.2)),
    ("Impact",       Inches(6.5),  Inches(1.2)),
    ("Mitigation",   Inches(7.8),  Inches(5.0)),
]:
    add_text(s, text, x, Inches(1.65), w, Inches(0.35), size=10, bold=True, color=WHITE)

risks = [
    ("Fuite données utilisateur",     "Sécurité",      ORANGE,  "Critique", RED,    "JWT httpOnly · HTTPS · bcrypt · Zod · Prisma ORM"),
    ("Rupture API Google News RSS",   "Disponibilité", RED,     "Modéré",   ORANGE, "Cache serveur · Sources RSS de secours · Erreur gracieuse"),
    ("Perte données PostgreSQL",      "Intégrité",     GREEN,   "Critique", RED,    "Backups · Docker volumes persistants · Prisma migrations"),
    ("Surcharge serveur (pics trafic)","Performance",  ORANGE,  "Modéré",   ORANGE, "Cache NHTSA 24h · Pagination API · Rate limiting prévu"),
    ("Non-conformité RGPD",           "Légal",         ORANGE,  "Élevé",    RED,    "Suppression compte · Minimisation données · Mentions légales"),
]
ry = Inches(2.05)
for risk, cat, cat_c, impact, imp_c, mitigation in risks:
    add_rect(s, Inches(0.22), ry, Inches(12.8), Inches(0.82),
             fill=DARK, line_color=RGBColor(0x33,0x33,0x33), line_width=Pt(0.5))
    add_text(s, risk,       Inches(0.35), ry+Inches(0.08), Inches(3.2),  Inches(0.65), size=11, bold=True,  color=WHITE)
    add_rect(s, Inches(3.65), ry+Inches(0.18), Inches(1.4),  Inches(0.42), fill=DARK,  line_color=LIGHTGRAY, line_width=Pt(0.5))
    add_text(s, cat,        Inches(3.68), ry+Inches(0.18), Inches(1.35), Inches(0.42), size=10, color=GRAY,  align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.22), ry+Inches(0.18), Inches(1.15), Inches(0.42), fill=cat_c)
    add_text(s, "Moyen" if cat_c==ORANGE else ("Élevé" if cat_c==RED else "Faible"),
             Inches(5.24), ry+Inches(0.18), Inches(1.1), Inches(0.42), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.5),  ry+Inches(0.18), Inches(1.15), Inches(0.42), fill=imp_c)
    add_text(s, impact, Inches(6.52), ry+Inches(0.18), Inches(1.1), Inches(0.42),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, mitigation, Inches(7.8), ry+Inches(0.08), Inches(5.1), Inches(0.65), size=10, color=GRAY)
    ry += Inches(0.88)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — C1.3.1 / C1.3.2 VEILLE & CHOIX TECHNO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.3.1 / C1.3.2 — VEILLE & CHOIX TECHNOLOGIQUES")
slide_title(s, "Pourquoi cette stack ?")

choices = [
    ("Next.js 14\nApp Router",    "vs Nuxt / Remix",
     "SSR + SSG natif, API Routes intégrées, écosystème React mature. App Router = layouts imbriqués, streaming.",
     "Tree-shaking avancé, bundle minimal, Edge Middleware natif"),
    ("PostgreSQL\n+ Prisma",      "vs MongoDB",
     "Données relationnelles (User → Moto → Entretien). Prisma = migrations versionnées, type-safe queries.",
     "Requêtes optimisées, indexes automatiques, ACID"),
    ("TypeScript\n+ Zod",         "vs JavaScript pur",
     "Typage statique = moins de bugs runtime. Zod valide les formulaires ET les réponses API avec le même schéma.",
     "Maintenance facilitée, refactoring sûr, autocomplétion"),
    ("shadcn/ui\n+ Tailwind",     "vs MUI / Bootstrap",
     "Composants injectés dans le projet (pas de dépendance externe). Tailwind = CSS minimal, pas de feuille globale.",
     "CSS purgé à 100% en prod, bundle JS réduit"),
]
positions = [
    (Inches(0.22), Inches(1.6)),
    (Inches(6.72), Inches(1.6)),
    (Inches(0.22), Inches(4.35)),
    (Inches(6.72), Inches(4.35)),
]
for (title, vs, desc, benefit), (x, y) in zip(choices, positions):
    card(s, x, y, Inches(6.3), Inches(2.6))
    # Title
    add_text(s, title, x+Inches(0.15), y+Inches(0.1), Inches(2.5), Inches(0.8),
             size=16, bold=True, color=RED)
    add_text(s, vs,    x+Inches(0.15), y+Inches(0.85), Inches(2.5), Inches(0.3),
             size=10, italic=True, color=LIGHTGRAY)
    # Desc
    bullet(s, [desc], x+Inches(0.15), y+Inches(1.15), Inches(6.0), Inches(0.9), size=11)
    # Green benefit
    add_rect(s, x+Inches(0.15), y+Inches(2.1), Inches(5.9), Inches(0.38),
             fill=RGBColor(0x0e,0x22,0x18), line_color=GREEN, line_width=Pt(0.5))
    add_text(s, "🌿  " + benefit, x+Inches(0.25), y+Inches(2.12), Inches(5.7), Inches(0.35),
             size=10, color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — C1.4.1 ESTIMATION DE CHARGE (CORRIGÉ : 29j)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.4.1 — FONCTIONNALITÉS & ESTIMATION DE CHARGE")
slide_title(s, "Périmètre fonctionnel")

cols_data = [
    ("PRINCIPALE",       RED,   [
        ("Authentification JWT (register/login/logout)", "3 j"),
        ("Gestion garage : wizard 3 étapes, liste, isPrimary", "5 j"),
        ("Historique maintenances (CRUD + 7 types)", "4 j"),
        ("Tableau de bord utilisateur (données réelles)", "3 j"),
    ]),
    ("SECONDAIRE",       BLUE,  [
        ("Page actualités (4 flux RSS Google News)", "2 j"),
        ("Recherche pièces (eBay/LBC/Amazon/Google)", "2 j"),
        ("Météo (API externe wttr.in)", "1 j"),
        ("Autocomplete modèles NHTSA (cache 24h)", "2 j"),
    ]),
    ("COMPLÉMENTAIRE",   GREEN, [
        ("Landing page 3D (Ducati V4R / Three.js)", "3 j"),
        ("Immatriculation auto-formatée SIV", "1 j"),
        ("27 marques + drapeaux (dont 12 chinoises)", "1 j"),
        ("Tests unitaires Jest + CI/CD GitHub Actions", "2 j"),
    ]),
]
cw = Inches(4.1)
cx = Inches(0.22)
for title, color, items in cols_data:
    card(s, cx, Inches(1.6), cw, Inches(5.1))
    card_header(s, title, cx, Inches(1.6), cw, fill=color)
    iy = Inches(2.1)
    for feat, days in items:
        # Day badge
        add_rect(s, cx+cw-Inches(0.65), iy, Inches(0.55), Inches(0.38), fill=color)
        add_text(s, days, cx+cw-Inches(0.65), iy+Inches(0.02), Inches(0.55), Inches(0.35),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, feat, cx+Inches(0.1), iy+Inches(0.04), cw-Inches(0.85), Inches(0.65),
                 size=10, color=GRAY)
        iy += Inches(1.1)
    cx += cw + Inches(0.13)

footer_note(s, "CHARGE TOTALE ESTIMÉE : 29 jours/homme  •  Soit ~6 semaines à temps plein  •  Équipe : 5 personnes")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — C1.4.2 BUDGET (CORRIGÉ : 29j × 350€)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.4.2 — BUDGET PRÉVISIONNEL")
slide_title(s, "Estimation financière du projet")

items = [
    ("Développement (29 j × 350 €)",                    10150, RED),
    ("Infrastructure cloud (Vercel + Supabase / an)",     1200, BLUE),
    ("Licences & outils (domaine, fonts, CDN)",            350, ORANGE),
    ("Tests & recette (3 j × 350 €)",                     1050, GREEN),
    ("Réserve risques (15%)",                             1900, LIGHTGRAY),
    ("Documentation & livraison (1 j)",                    350, RGBColor(0x88,0x44,0xAA)),
]
total = sum(v for _, v, _ in items)
max_v = max(v for _, v, _ in items)

by = Inches(1.6)
for label, value, color in items:
    add_text(s, label, Inches(0.35), by+Inches(0.08), Inches(6.5), Inches(0.5), size=12, color=GRAY)
    bar_w = Inches(0.3) + Inches(5.5) * (value / max_v)
    add_rect(s, Inches(7.0), by+Inches(0.1), bar_w, Inches(0.38), fill=color)
    add_text(s, f"{value:,} €".replace(",", " "), Inches(7.0)+bar_w+Inches(0.1),
             by+Inches(0.1), Inches(1.5), Inches(0.38), size=12, bold=True, color=WHITE)
    by += Inches(0.75)

# Total
add_rect(s, Inches(0.22), Inches(6.6), Inches(12.8), Inches(0.65),
         fill=RGBColor(0x18,0x18,0x18), line_color=RED, line_width=Pt(1.5))
add_text(s, "BUDGET TOTAL PRÉVISIONNEL", Inches(0.35), Inches(6.65), Inches(7), Inches(0.5),
         size=13, bold=True, color=WHITE)
add_text(s, f"{total:,} €".replace(",", " "), Inches(9.0), Inches(6.6), Inches(3.8), Inches(0.65),
         size=32, bold=True, color=RED, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — C1.6 PRÉCONISATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); red_bar(s)
top_label(s, "C1.6 — PRÉCONISATION & ARGUMENTATION")
slide_title(s, "Recommandations & prochaines étapes")

phases = [
    ("Phase 1 — 0 à 2 mois",  RED,  [
        "Déploiement Vercel + Supabase (PostgreSQL managé)",
        "HTTPS + headers sécurité (CSP, HSTS, X-Frame)",
        "Tests E2E Playwright sur parcours critiques",
        "Conformité RGPD — mentions légales + suppression compte",
    ]),
    ("Phase 2 — 2 à 4 mois",  BLUE, [
        "OAuth Google/Apple en complément JWT",
        "Notifications push rappels maintenance (PWA)",
        "Mode hors-ligne garage (Service Worker)",
        "Export PDF carnet d'entretien",
    ]),
    ("Phase 3 — 4 à 6 mois",  GREEN,[
        "Application mobile React Native (partage de code)",
        "Monétisation : plan premium (multi-garage, export)",
        "Analytics usage anonymisé (Plausible.io)",
        "Recherche full-text PostgreSQL sur les entretiens",
    ]),
]
cw = Inches(4.1)
cx = Inches(0.22)
for title, color, items in phases:
    card(s, cx, Inches(1.6), cw, Inches(5.1))
    card_header(s, title, cx, Inches(1.6), cw, fill=color)
    bullet(s, items, cx+Inches(0.15), Inches(2.12), cw-Inches(0.3), Inches(4.4), size=11.5)
    cx += cw + Inches(0.13)

footer_note(s, "MotoTrack est une solution viable, différenciante et maintenable. Next.js + Prisma + PostgreSQL offre un rapport qualité/coût/performance optimal.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION (CORRIGÉ)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
add_rect(s, 0, 0, Inches(0.12), H, fill=RED)
add_rect(s, 0, H-Inches(0.12), W, Inches(0.12), fill=RED)

add_text(s, "MotoTrack", Inches(0.3), Inches(0.8), W-Inches(0.5), Inches(2.0),
         size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Une application pensée pour les passionnés de moto",
         Inches(0.3), Inches(2.7), W-Inches(0.5), Inches(0.5),
         size=20, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

stats = [
    ("8",    "Pages Next.js"),
    ("9",    "API endpoints"),
    ("29 j", "Charge estimée"),
    ("15k€", "Budget prévu"),
]
sx = Inches(0.8)
for val, label in stats:
    add_rect(s, sx, Inches(3.4), Inches(2.7), Inches(1.6),
             fill=DARK, line_color=RGBColor(0x44,0x44,0x44))
    add_text(s, val,   sx, Inches(3.45), Inches(2.7), Inches(0.9),
             size=48, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, label, sx, Inches(4.25), Inches(2.7), Inches(0.55),
             size=12, color=GRAY, align=PP_ALIGN.CENTER)
    sx += Inches(2.95)

add_text(s, "Merci pour votre attention — Questions ?",
         Inches(0.3), Inches(5.3), W-Inches(0.5), Inches(0.6),
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.0), Inches(6.0), Inches(3.3), Inches(0.06), fill=RED)

# ─── Save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\lucas\OneDrive\Documents\Ynov_M2\MotoTrack\MotoTrack_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
